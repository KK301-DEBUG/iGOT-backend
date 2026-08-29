"""
Text extraction from uploaded training materials.
Supports PDF, DOCX, and PPTX.
"""

import fitz
import docx
import pptx

# extracting text from pdf
def extract_from_pdf(path) -> str:
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

# extracting text from docx
def extract_from_docx(path: str) -> str:
    document = docx.Document(path)
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())

# extracting text from pdf
def extract_from_pptx(path: str) -> str:
    ppt = pptx.Presentation(path)
    text_runs = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        text_runs.append(line)
    return "\n".join(text_runs)

# dispatch extraction based on file extension
def extract_text(path: str) -> str:
    ext = path.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_from_pdf(path)
    elif ext == "docx":
        return extract_from_docx(path)
    elif ext == "pptx":
        return extract_from_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")